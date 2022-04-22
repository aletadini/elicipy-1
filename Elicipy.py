import datetime
import copy
from pptx.parts.chart import ChartPart
from pptx.parts.embeddedpackage import EmbeddedXlsxPart
from difflib import SequenceMatcher
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx import Presentation
import difflib
from scipy import stats
from matplotlib.ticker import PercentFormatter
from script_fromR import createDATA1
from global_weights import global_weights
from script_fromR import generate_ERF
from merge_csv import merge_csv
import numpy as np
import pkg_resources
import re
import os
import sys
from pdf2image import convert_from_path
from pathlib import Path
import openpyxl
import pandas as pd
import matplotlib
import matplotlib.pyplot as plt

from ElicipyDict import *
from matplotlib import rcParams

rcParams['font.family'] = 'sans-serif'
rcParams['font.sans-serif'] = ['Helvetica']

matplotlib.use("TkAgg")

def add_date(slide):

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.2),
                                           Inches(16), Inches(0.3))
    shape.shadow.inherit = False
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(60, 90, 180)
    line = shape.line
    line.color.rgb = RGBColor(60, 90, 180)
    shape.text = "Expert elicitation " + \
        datetime.datetime.today().strftime('%d-%b-%Y')
    shape_para = shape.text_frame.paragraphs[0]
    shape_para.font.name = "Helvetica"
    shape_para.font.size = Pt(17)

def add_small_logo(slide,left,top):

    img = slide.shapes.add_picture('logo.png',left + Inches(13.0),
                                   top+Inches(6.8),
                                   width=Inches(0.8))

def add_figure(slide,figname,left,top):

    img = slide.shapes.add_picture(figname,left + Inches(3.4),
                                           top,
                                           width=Inches(10))

def add_title(slide,text_title):

    title_shape = slide.shapes.title
    title_shape.text = text_title
    title_shape.width = Inches(15)
    title_shape.height = Inches(2)
    title_para = slide.shapes.title.text_frame.paragraphs[0]
    title_para.font.name = "Helvetica"

def add_text_box(slide,left,top,text_box):

    txBox = slide.shapes.add_textbox(left - Inches(1),
                                     top + Inches(0.5),
                                     width=Inches(4),
                                     height=Inches(5))
    tf = txBox.text_frame

    p = tf.add_paragraph()
    p.text = text_box
    p.font.name = "Helvetica"
    p.font.size = Pt(16)


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


# get current path
path = os.getcwd()

os.chdir(path)

# change to full path
output_dir = path + '/' + output_dir
input_dir = path + '/' + input_dir

# merge the files of the different experts
# creating one file for seeds and one for tagets
merge_csv(input_dir, target)

# Check whether the specified output path exists or not
isExist = os.path.exists(output_dir)

if not isExist:

    # Create a new directory because it does not exist
    os.makedirs(output_dir)
    print('The new directory ' + output_dir + ' is created!')

# seeds file name
filename = input_dir + '/seed.csv'

# Read a comma-separated values (csv) file into DataFrame df_SQ
df_SQ = pd.read_csv(filename)

# The second column (index 1) must contain the first name of the expert
firstname = df_SQ[df_SQ.columns[1]].astype(str).tolist()

# The third column (index 2) must contain the first name of the expert
surname = df_SQ[df_SQ.columns[2]].astype(str).tolist()

# create a list with firstname+surname
# this is needed to search for expert matches between seed and target questions
NS_SQ = []

for name, surname in zip(firstname, surname):

    NS_SQ.append(name + surname)

print('NS_SQ', NS_SQ)

# create a 2D numpy array with the answers to the seed questions
cols_as_np = df_SQ[df_SQ.columns[3:]].to_numpy()

# we want to work with a 3D array, with the following dimension:
# n_expert X n_pctl X n_SQ
n_experts = cols_as_np.shape[0]
n_pctl = 3
n_SQ = int(cols_as_np.shape[1] / n_pctl)

# reshaped numpy array with expert answers
SQ_array = np.reshape(cols_as_np, (n_experts, n_SQ, n_pctl))

# swap the array to have seed for the last index
SQ_array = np.swapaxes(SQ_array, 1, 2)

# sort according to the percentile values
# (sometimes the expert give the percentiles in the wrong order)
SQ_array = np.sort(SQ_array, axis=1)

df_quest = pd.read_csv(input_dir + '/' + csv_file, header=0)

# list with the short title of the target questions
SQ_question = []
# list with the long title of the target questions
SQ_LongQuestion = []
# list with min vals for target questions
SQ_minVals = []
# list with max vals for target questions
SQ_maxVals = []
# list with the units of the target questions
SQ_units = []
# scale for target question:
SQ_scale = []

SQ_realization = []

global_log = []

for i in df_quest.itertuples():

    idx, shortQ, longQ, unit, scale, minVal, maxVal, realization, question = i[
        0:9]

    if scale == 'uni':

        global_log.append(0)

    else:

        global_log.append(1)

    if (question == 'seed'):

        SQ_question.append(shortQ)
        SQ_LongQuestion.append(longQ)
        SQ_units.append(unit)
        SQ_scale.append(scale)

        SQ_realization.append(realization)

        if minVal.is_integer():

            minVal = int(minVal)

        if maxVal.is_integer():

            maxVal = int(maxVal)

        SQ_minVals.append(minVal)
        SQ_maxVals.append(maxVal)

# print on screen the units
print("Seed_units = ", SQ_units)

# print on screen the units
print("Seed_scales = ", SQ_scale)

for i in np.arange(n_SQ):

    for k in np.arange(n_experts):

        # if 5% and 50% percentiles are equal, reduce 5%
        if SQ_array[k, 0, i] == SQ_array[k, 1, i]:

            SQ_array[k, 0, i] = SQ_array[k, 1, i] * 0.99

        # if 50% and 95% percentiles are equal, increase 95%
        if SQ_array[k, 2, i] == SQ_array[k, 1, i]:

            SQ_array[k, 2, i] = SQ_array[k, 1, i] * 1.01

    print('')
    print('Seed question ', i)
    print(SQ_array[:, :, i])
    # print('mean(lin) =',np.nanmean(check_lin),'mean(log) =',np.nanmean(check_log))

# list with the "title" of the target questions
TQ_question = []
# list with the long title of the target questions
TQ_LongQuestion = []
TQ_minVals = []
TQ_maxVals = []
# list with the units of the target questions
TQ_units = []
# scale for target question:
TQ_scale = []

if target:

    filename = input_dir + '/target.csv'

    # Read a comma-separated values (csv) file into DataFrame df_TQ
    df_TQ = pd.read_csv(filename)

    # The second column (index 1) must contain the first name of the expert
    firstname = df_TQ[df_TQ.columns[1]].astype(str).tolist()

    # The third column (index 2) must contain the first name of the expert
    surname = df_TQ[df_TQ.columns[2]].astype(str).tolist()

    # create a list with firstname+surname
    # this is needed to search for expert matches between seed and target questions
    NS_TQ = []

    for name, surname in zip(firstname, surname):

        NS_TQ.append(name + surname)

    print('NS_TQ', NS_TQ)

    sorted_idx = []

    # loop to search for matches between experts in seed and target
    for TQ_name in NS_TQ:

        index = NS_SQ.index(difflib.get_close_matches(TQ_name, NS_SQ)[0])
        sorted_idx.append(index)

    print('Sorted list of experts to match the order of seeds:', sorted_idx)

    # create a 2D numpy array with the answers to the seed questions
    cols_as_np = df_TQ[df_TQ.columns[3:]].to_numpy()

    # sort for expert names
    cols_as_np = cols_as_np[sorted_idx, :]

    # we want to work with a 3D array, with the following dimension:
    # n_expert X n_pctl X n_TQ
    n_experts_TQ = cols_as_np.shape[0]

    # check if number of experts in seed and target is the same
    if (n_experts_TQ != n_experts):

        print('Error: number of experts in seeds and targets different')
        sys.exit()

    n_TQ = int(cols_as_np.shape[1] / n_pctl)

    # reshaped numpy array with expert answers
    TQ_array = np.reshape(cols_as_np, (n_experts, n_TQ, n_pctl))

    # swap the array to have pctls for the second index
    TQ_array = np.swapaxes(TQ_array, 1, 2)

    # sort according to the percentile values
    # (sometimes the expert give the percentiles in the wrong order)
    TQ_array = np.sort(TQ_array, axis=1)

    for i in df_quest.itertuples():

        idx, shortQ, longQ, unit, scale, minVal, maxVal, realization, question = i[
            0:9]

        if (question == 'target'):

            TQ_question.append(shortQ)
            TQ_LongQuestion.append(longQ)
            TQ_units.append(unit)
            TQ_scale.append(scale)

            if minVal.is_integer():

                minVal = int(minVal)

            if maxVal.is_integer():

                maxVal = int(maxVal)

            TQ_minVals.append(minVal)
            TQ_maxVals.append(maxVal)

    # print on screen the units
    print("Target units = ", TQ_units)

    # print on screen the units
    print("Target scales = ", TQ_scale)

    global_scale = SQ_scale + TQ_scale

    for i in np.arange(n_TQ):

        for k in np.arange(n_experts):
            if TQ_array[k, 0, i] == TQ_array[k, 1, i]:
                TQ_array[k, 0, i] = TQ_array[k, 1, i] * 0.99
            if TQ_array[k, 2, i] == TQ_array[k, 1, i]:
                TQ_array[k, 2, i] = TQ_array[k, 1, i] * 1.01

        print('Target question ', i)
        print(TQ_array[:, :, i])

else:

    # if we do not read the target questions, set empty array
    n_TQ = 0

    TQ_array = np.zeros((n_experts, n_pctl, n_TQ))
    global_scale = SQ_scale

if target:

    nTot = TQ_array.shape[2] + SQ_array.shape[2]

else:

    nTot = SQ_array.shape[2]

realization = np.zeros(TQ_array.shape[2] + SQ_array.shape[2])
realization[0:SQ_array.shape[2]] = SQ_realization

global_minVal = SQ_minVals + TQ_minVals
global_maxVal = SQ_maxVals + TQ_maxVals
global_units = SQ_units + TQ_units
global_longQuestion = SQ_LongQuestion + TQ_LongQuestion
global_shortQuestion = SQ_question + TQ_question

print("")
print('Realization', realization)

# ----------------------------------------- #
# ------------ Compute weights ------------ #
# ----------------------------------------- #


if analysis:

    W = global_weights(SQ_array, TQ_array, realization, alpha, global_scale, k,
                       cal_power)

    W_erf = generate_ERF(realization, SQ_array)

    Weq = np.ones(n_experts)
    Weqok = [x / n_experts for x in Weq]

    W_gt0_01 = []
    expin = []

    for x in W[:, 4]:
        if x > 0:
            W_gt0_01.append(x)

    k = 1
    for i in W[:, 4]:
        if i > 0:
            expin.append(k)
        k += 1

    W_gt0 = [round((x * 100), 1) for x in W_gt0_01]

    Werf_gt0_01 = []
    expin = []

    for x in W_erf[:, 4]:
        if x > 0:
            Werf_gt0_01.append(x)

    k = 1
    for i in W_erf[:, 4]:
        if i > 0:
            expin.append(k)
        k += 1

    Werf_gt0 = [round((x * 100), 1) for x in Werf_gt0_01]

    print("")
    print('W_erf')
    print(W_erf[:, -1])
    print("")
    print('W')
    print(W[:, -1])
    print("")
    print('Weq')
    print(Weqok)

# ----------------------------------------- #
# ------ Create samples and bar plots ----- #
# ----------------------------------------- #

DAT = np.zeros((n_experts * (n_SQ + n_TQ), n_pctl + 2))

DAT[:, 0] = np.repeat(np.arange(1, n_experts + 1), n_SQ + n_TQ)
DAT[:, 1] = np.tile(np.arange(1, n_SQ + n_TQ + 1), n_experts)

DAT[:, 2:] = np.append(SQ_array, TQ_array, axis=2).transpose(0, 2,
                                                             1).reshape(-1, 3)

q05 = []
q50 = []
q95 = []

q05_erf = []
q50_erf = []
q95_erf = []

q05_EW = []
q50_EW = []
q95_EW = []

figs_h = {}
axs_h = {}
axs_h2 = {}

plt.rcParams.update({'font.size': 8})

samples = np.zeros((n_sample, n_TQ))
samples_erf = np.zeros((n_sample, n_TQ))
samples_EW = np.zeros((n_sample, n_TQ))

print("")
if analysis:
    print(" j   quan05    quan50     qmean    quan95")

del_rows = []
keep_rows = []

if (not Cooke_flag):
    del_rows.append(int(0))
else:
    keep_rows.append(int(0))
if (not ERF_flag):
    del_rows.append(int(1))
else:
    keep_rows.append(int(1))
if (not EW_flag):
    del_rows.append(int(2))
else:
    keep_rows.append(int(2))

colors = ['orange', 'purple', 'springgreen']
colors = [colors[index] for index in keep_rows]

legends = ['CM', 'ERF', 'EW']
legends = [legends[index] for index in keep_rows]

for j in np.arange(n_SQ + n_TQ):

    if analysis:

        quan05, quan50, qmean, quan95, C = createDATA1(
            DAT, j, W[:, 4].flatten(), n_sample, global_log[j],
            [global_minVal[j], global_maxVal[j]], False)

        print("%2i %9.2f %9.2f %9.2f %9.2f" %
              (j, quan05, quan50, qmean, quan95))

        q05.append(quan05)
        q50.append(quan50)
        q95.append(quan95)

        quan05_erf, quan50_erf, qmean_erf, quan95_erf, C_erf = createDATA1(
            DAT, j, W_erf[:, 4].flatten(), n_sample, global_log[j],
            [global_minVal[j], global_maxVal[j]], True)

        print("%2i %9.2f %9.2f %9.2f %9.2f" %
              (j, quan05_erf, quan50_erf, qmean_erf, quan95_erf))

        q05_erf.append(quan05_erf)
        q50_erf.append(quan50_erf)
        q95_erf.append(quan95_erf)

        quan05_EW, quan50_EW, qmean_EW, quan95_EW, C_EW = createDATA1(
            DAT, j, Weqok, n_sample, global_log[j],
            [global_minVal[j], global_maxVal[j]], False)

        print("%2i %9.2f %9.2f %9.2f %9.2f" %
              (j, quan05_EW, quan50_EW, qmean_EW, quan95_EW))

        q05_EW.append(quan05_EW)
        q50_EW.append(quan50_EW)
        q95_EW.append(quan95_EW)

        if (j >= n_SQ):

            samples[:, j - n_SQ] = C
            samples_erf[:, j - n_SQ] = C_erf
            samples_EW[:, j - n_SQ] = C_EW

        if (j >= n_SQ):

            ntarget = str(j - n_SQ + 1)

            figs_h[j] = plt.figure()
            axs_h[j] = figs_h[j].add_subplot(111)
            C_stack = np.stack((C, C_erf, C_EW), axis=0)
            C_stack = np.delete(C_stack, del_rows, 0)
            wg = np.ones_like(C_stack.T) / n_sample

            if hist_type == 'step':

                axs_h[j].hist(C_stack.T,
                              bins=n_bins,
                              weights=wg,
                              histtype='step',
                              fill=False,
                              rwidth=0.95,
                              color=colors)

            elif hist_type == 'bar':

                axs_h[j].hist(C_stack.T,
                              bins=n_bins,
                              weights=wg,
                              histtype='bar',
                              rwidth=0.95,
                              ec="k",
                              color=colors)

            axs_h[j].set_xlabel(TQ_units[j - n_SQ])
            plt.gca().yaxis.set_major_formatter(PercentFormatter(1))

            xt = plt.xticks()[0]

            axs_h2[j] = axs_h[j].twinx()

            if global_units[j] == "%":

                xmin = 0.0
                xmax = 100.0

            else:

                xmin = np.amin(C_stack)
                xmax = np.amax(C_stack)

            lnspc = np.linspace(xmin, xmax, 1000)

            if (Cooke_flag):
                gkde = stats.gaussian_kde(C)
                gkde_norm = gkde.integrate_box_1d(global_minVal[j],
                                                  global_maxVal[j])
                kdepdf = gkde.evaluate(lnspc) / gkde_norm
                axs_h2[j].plot(lnspc, kdepdf, 'r--')

            if (ERF_flag):
                gkde_erf = stats.gaussian_kde(C_erf)
                gkde_erf_norm = gkde_erf.integrate_box_1d(
                    global_minVal[j], global_maxVal[j])
                kdepdf_erf = gkde_erf.evaluate(lnspc) / gkde_erf_norm
                axs_h2[j].plot(lnspc, kdepdf_erf, '--', color='tab:purple')

            if (EW_flag):
                gkde_EW = stats.gaussian_kde(C_EW)
                gkde_EW_norm = gkde_EW.integrate_box_1d(
                    global_minVal[j], global_maxVal[j])
                kdepdf_EW = gkde_EW.evaluate(lnspc) / gkde_EW_norm
                axs_h2[j].plot(lnspc, kdepdf_EW, 'g--')

            axs_h[j].set_xlim(xmin, xmax)
            axs_h2[j].set_xlim(xmin, xmax)

            axs_h2[j].set_ylabel('PDF', color='b')

            axs_h2[j].set_ylim(bottom=0)
            plt.legend(legends)
            plt.title('Target Question ' + str(j - n_SQ + 1))

            figname = output_dir + '/' + elicitation_name + \
                '_hist_' + str(j - n_SQ + 1).zfill(2) + '.pdf'
            figs_h[j].savefig(figname)

            images = convert_from_path(figname)
            figname = output_dir + '/' + elicitation_name + \
                '_hist_' + str(j - n_SQ + 1).zfill(2) + '.png'
            images[0].save(figname, 'PNG')

            plt.close()

# ----------------------------------------- #
# ---------- Save samples on csv ---------- #
# ----------------------------------------- #

if Cooke_flag:

    csv_name = output_dir + '/' + elicitation_name + '_samples.csv'
    np.savetxt(csv_name, samples, delimiter=",", fmt='%1.4e')

if ERF_flag:

    csv_name = output_dir + '/' + elicitation_name + '_samples_erf.csv'
    np.savetxt(csv_name, samples_erf, delimiter=",", fmt='%1.4e')

if EW_flag:

    csv_name = output_dir + '/' + elicitation_name + '_samples_EW.csv'
    np.savetxt(csv_name, samples_EW, delimiter=",", fmt='%1.4e')

# ----------------------------------------- #
# --------- Create answ. figures ---------- #
# ----------------------------------------- #

figs = {}
axs = {}

for h in np.arange(n_SQ + n_TQ):

    if (h >= n_SQ):

        j = h - n_SQ
        Q_array = TQ_array[:, :, j]
        string = 'Target'

    else:

        j = h
        Q_array = SQ_array[:, :, j]
        string = 'Seed'

    x = Q_array[:, 1]
    y = np.arange(n_experts) + 1

    # creating error
    x_errormax = Q_array[:, 2] - Q_array[:, 1]
    x_errormin = Q_array[:, 1] - Q_array[:, 0]

    x_error = [x_errormin, x_errormax]

    figs[j] = plt.figure()
    axs[j] = figs[j].add_subplot(111)
    axs[j].errorbar(x, y, xerr=x_error, fmt='bo')
    axs[j].plot(x - x_errormin, y, 'bx')
    axs[j].plot(x + x_errormax, y, 'bx')

    if (realization[j] > 999):
        txt = '%5.2e' % realization[h]
    else:
        txt = '%6.2f' % realization[h]

    ytick = []
    for i in y:
        ytick.append('Exp.' + str(int(i)))

    yerror = n_experts
    if analysis:

        if Cooke_flag:

            yerror = yerror + 1
            axs[j].errorbar(q50[h],
                            yerror,
                            xerr=[[q50[h] - q05[h]], [q95[h] - q50[h]]],
                            fmt='ro')
            axs[j].plot(q05[h], yerror, 'rx')
            axs[j].plot(q95[h], yerror, 'rx')

            ytick.append('DM-Cooke')

        if ERF_flag:

            yerror = yerror + 1
            axs[j].errorbar(q50_erf[h], [yerror],
                            xerr=[[q50_erf[h] - q05_erf[h]],
                                  [q95_erf[h] - q50_erf[h]]],
                            fmt='o',
                            color='tab:purple')
            axs[j].plot(q05_erf[h], yerror, 'x', color='tab:purple')
            axs[j].plot(q95_erf[h], yerror, 'x', color='tab:purple')

            ytick.append('DM-ERF')

        if EW_flag:

            yerror = yerror + 1
            axs[j].errorbar([q50_EW[h]], [yerror],
                            xerr=[[q50_EW[h] - q05_EW[h]],
                                  [q95_EW[h] - q50_EW[h]]],
                            fmt='go')

            axs[j].plot(q05_EW[h], yerror, 'gx')
            axs[j].plot(q95_EW[h], yerror, 'gx')

            ytick.append('DM-Equal')

        if (h < n_SQ):

            yerror = yerror + 1
            axs[j].plot(realization[h], yerror, 'kx')
            axs[j].annotate(txt, (realization[h], yerror + 0.15))

            ytick.append('Realization')

    else:

        if (h < n_SQ):

            axs[j].plot(realization[h], n_experts + 1, 'kx')
            axs[j].annotate(txt, (realization[j], yerror + 0.15))

    y = np.arange(len(ytick)) + 1

    ytick_tuple = tuple(i for i in ytick)
    axs[j].set_yticks(y)

    axs[j].set_yticklabels(ytick_tuple)
    axs[j].set_xlabel(global_units[h])

    if (global_log[h] == 1):

        axs[j].set_xscale('log')

    axs[j].set_ylim(0.5, len(ytick) + 0.5)

    axs[j].grid(linewidth=0.4)

    plt.title(string + ' Question ' + str(j + 1))
    figname = output_dir + '/' + elicitation_name + \
        '_'+string+'_' + str(j + 1).zfill(2) + '.pdf'
    figs[j].savefig(figname)

    images = convert_from_path(figname)
    figname = output_dir + '/' + elicitation_name + \
        '_'+string+'_' + str(j + 1).zfill(2) + '.png'
    images[0].save(figname, 'PNG')
    plt.close()

# ----------------------------------------- #
# ------- Create .pptx presentation ------- #
# ----------------------------------------- #

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
left = Inches(2)
top = Inches(1.5)

# ------------- Title slide ----------------#
lyt = prs.slide_layouts[0]  # choosing a slide layout
slide = prs.slides.add_slide(lyt)  # adding a slide
title = slide.shapes.title  # assigning a title
subtitle = slide.placeholders[1]  # placeholder for subtitle
title.text = "Expert elicitation - "+elicitation_name  # title

title_para = slide.shapes.title.text_frame.paragraphs[0]
title_para.font.name = "Helvetica"

Current_Date_Formatted = datetime.datetime.today().strftime('%d-%b-%Y')

subtitle.text = Current_Date_Formatted  # subtitle

subtitle_para = slide.shapes.placeholders[1].text_frame.paragraphs[0]
subtitle_para.font.name = "Helvetica"

img = slide.shapes.add_picture('logo.png',left + Inches(11.3),
                                   top+Inches(5.4),
                                   width=Inches(2.4))

title_slide_layout = prs.slide_layouts[5]

# ------------- Weights slide -------------#

if analysis:

    slide = prs.slides.add_slide(title_slide_layout)
    
    text_title = "Experts' weights"
    add_title(slide,text_title)
    
    # ---add table weights to slide---
    x, y, cx, cy = Inches(2), Inches(2), Inches(8), Inches(4)

    shape = slide.shapes.add_table(3,
                                   len(W_gt0) + 1, x, y, cx,
                                   MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)

    table = shape.table

    cell = table.cell(0, 0)
    cell.text = 'Expert ID'

    cell = table.cell(1, 0)
    cell.text = 'Expert weight (Cooke)'

    cell = table.cell(2, 0)
    cell.text = 'Expert weight (ERF)'

    for j in np.arange(len(W_gt0)):
        cell = table.cell(0, j + 1)
        cell.text = 'Exp' + str(expin[j])

        cell = table.cell(1, j + 1)
        cell.text = '%6.2f' % W_gt0[j]

        cell = table.cell(2, j + 1)
        cell.text = '%6.2f' % Werf_gt0[j]

    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)

    add_date(slide)            
    add_small_logo(slide,left,top)

# ------------- Answers slides ------------#

for h in np.arange(n_SQ + n_TQ):

    if (h >= n_SQ):

        j = h - n_SQ
        string = 'Target'

    else:

        j = h
        string = 'Seed'

    slide = prs.slides.add_slide(title_slide_layout)

    text_title = global_shortQuestion[h]
    add_title(slide,text_title)

    text_box = global_longQuestion[h]
    add_text_box(slide,left,top,text_box)

    figname = output_dir + '/' + elicitation_name + \
        '_'+string+'_' + str(j + 1).zfill(2) + '.png'
    add_figure(slide,figname,left,top)

    add_date(slide)            
    add_small_logo(slide,left,top)

# ------------- Pctls slides -------------#

if analysis and target:

    slide = prs.slides.add_slide(prs.slide_layouts[5])
        
    text_title = "Percentiles of target questions"
    add_title(slide,text_title)
    
    
    # ---add table to slide---
    x, y, cx, cy = Inches(2), Inches(2), Inches(12), Inches(4)
    shape = slide.shapes.add_table(n_TQ + 1, 7, x, y, cx,
                                   MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    table = shape.table

    cell = table.cell(0, 1)
    cell.text = 'Q05 (Cooke)'

    cell = table.cell(0, 2)
    cell.text = 'Q50 (Cooke)'

    cell = table.cell(0, 3)
    cell.text = 'Q95 (Cooke)'

    cell = table.cell(0, 4)
    cell.text = 'Q05 (ERF)'

    cell = table.cell(0, 5)
    cell.text = 'Q50 (ERF)'

    cell = table.cell(0, 6)
    cell.text = 'Q95 (ERF)'

    for h in np.arange(n_TQ):

        j = h + n_SQ

        cell = table.cell(h + 1, 0)
        cell.text = 'Target Question ' + str(h + 1)

        cell = table.cell(h + 1, 1)
        cell.text = '%6.2f' % q05[j]

        cell = table.cell(h + 1, 2)
        cell.text = '%6.2f' % q50[j]

        cell = table.cell(h + 1, 3)
        cell.text = '%6.2f' % q95[j]

        cell = table.cell(h + 1, 4)
        cell.text = '%6.2f' % q05_erf[j]

        cell = table.cell(h + 1, 5)
        cell.text = '%6.2f' % q50_erf[j]

        cell = table.cell(h + 1, 6)
        cell.text = '%6.2f' % q95_erf[j]

    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)
     
    add_date(slide)            
    add_small_logo(slide,left,top)

# ------------ Barplot slides ------------#

for j in np.arange(n_SQ + n_TQ):

    if analysis:

        if (j >= n_SQ):

            slide = prs.slides.add_slide(title_slide_layout)

            figname = output_dir + '/' + elicitation_name + \
                '_hist_' + str(j - n_SQ + 1).zfill(2) + '.png'
            
            text_title = TQ_question[j - n_SQ]
            add_title(slide,text_title)

            text_box = TQ_LongQuestion[j - n_SQ]
            add_text_box(slide,left,top,text_box)

            add_date(slide)
            add_small_logo(slide,left,top)
            add_figure(slide,figname,left-Inches(0.8),top)


prs.save(output_dir + "/" + elicitation_name + ".pptx")  # saving file
